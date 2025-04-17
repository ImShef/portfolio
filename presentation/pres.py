from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import io
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ========== СТИЛИ ==========
COLORS = {
    'primary': RGBColor(13, 59, 102),
    'secondary': RGBColor(87, 171, 39),
    'accent': RGBColor(255, 196, 37),
    'background': RGBColor(241, 245, 249)
}

FONTS = {
    'title': ('Calibri Light', Pt(36)),
    'body': ('Calibri', Pt(18))
}

# ========== ФУНКЦИИ ==========
def set_background(slide):
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = COLORS['background']
    fill.gradient_stops[1].color.rgb = RGBColor(255, 255, 255)

def add_title(slide, text):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    title.text_frame.paragraphs[0].font.size = FONTS['title'][1]
    title.text_frame.paragraphs[0].font.name = FONTS['title'][0]

def add_decorative_line(slide, start_x, start_y, length):
    line = slide.shapes.add_shape(
        MSO_SHAPE.LINE_INVERSE,
        left=Cm(start_x),
        top=Cm(start_y),
        width=Cm(length),
        height=Cm(0))
    line.line.color.rgb = COLORS['accent']
    line.line.width = Pt(2.5)

def add_content_box(slide, text, pos_y):
    textbox = slide.shapes.add_textbox(Cm(1.5), Cm(pos_y), Cm(24), Cm(4.5))
    tf = textbox.text_frame
    tf.word_wrap = True
    
    for line in text.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['primary']
        p.space_after = Pt(6)
    
    add_decorative_line(slide, 1.5, pos_y-0.3, 24)

# ========== СЛАЙД 1: Титульный ==========
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_background(slide)
add_title(slide, "Шумовое загрязнение:\nВлияние на экосистемы")

subtitle = slide.shapes.add_textbox(Cm(1.5), Cm(8), Cm(20), Cm(2))
subtitle.text_frame.text = "Экологическое исследование • 2023"
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
subtitle.text_frame.paragraphs[0].font.size = Pt(18)

# ========== СЛАЙД 2: Уровни шума ==========
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_background(slide)
add_title(slide, "Уровни звукового давления")

fig, ax = plt.subplots(figsize=(8,4))
labels = ['Шёпот (30 дБ)', 'Офис (60 дБ)', 'Поезд (90 дБ)']
values = [30, 60, 90]
ax.barh(labels, values, color=['#4a86e8', '#87ab39', '#ffc425'])
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.tick_params(colors='#0D3B66')
plt.close()

# ========== СЛАЙД 3: Воздействие на здоровье ==========
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_background(slide)
add_title(slide, "Физиологические эффекты")
content = """• Нарушения сна → снижение когнитивных функций
• Повышение уровня кортизола на 37%
• Риск гипертонии увеличивается в 1.5 раза
• Потеря слуха при длительном воздействии"""
add_content_box(slide, content, 4.5)

# ========== СЛАЙД 4: Источники шума ==========
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_background(slide)
add_title(slide, "Основные источники шума")
content = """🏭 Промышленность:
• Заводы
• Строительные площадки
• Горнодобывающие комплексы

🚗 Транспорт:
• Автомагистрали
• Аэропорты
• Железные дороги

🏙️ Городская среда:
• Торговые центры
• Концертные площадки
• Общепит"""
add_content_box(slide, content, 4)

# ========== СЛАЙД 5: Влияние на животных ==========
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_background(slide)
add_title(slide, "Экологические последствия")

fig, ax = plt.subplots(figsize=(6,6))
labels = ['Морские млекопитающие', 'Птицы', 'Насекомые']
sizes = [45, 30, 25]
ax.pie(sizes, labels=labels, autopct='%1.1f%%', colors=['#4a86e8', '#87ab39', '#ffc425'])



content = """🐳 Киты: потеря ориентации
🦉 Совы: снижение охотничьей активности
🐝 Пчелы: нарушение навигации
🐸 Лягушки: сбой репродукции"""
add_content_box(slide, content, 4)

# ========== СЛАЙД 6: Нормативы ==========
slide = prs.slides.add_slide(prs.slide_layouts[5])
set_background(slide)
add_title(slide, "Правовое регулирование")

table = slide.shapes.add_table(rows=4, cols=3, left=Cm(3), top=Cm(4), width=Cm(20), height=Cm(6)).table
table_data = [
    ['Зона', 'День (дБ)', 'Ночь (дБ)'],
    ['Жилая', '55', '45'],
    ['Промышленная', '75', '65'],
    ['Транспортная', '85', '75']
]

for row_idx, row in enumerate(table_data):
    for col_idx, cell in enumerate(row):
        table.cell(row_idx, col_idx).text = cell
        tf = table.cell(row_idx, col_idx).text_frame
        tf.paragraphs[0].font.color.rgb = COLORS['primary']
        tf.paragraphs[0].font.size = Pt(16)
        if row_idx == 0:
            table.cell(row_idx, col_idx).fill.solid()
            table.cell(row_idx, col_idx).fill.fore_color.rgb = COLORS['accent']


# ======== СЛАЙД 7: Технические решения ========
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_title(slide, "Технологии снижения шума")

content = [
    ["🌉 Шумовые экраны", 
     "• Поглощение до 15 дБ\n• Композитные материалы\n• Зеленые насаждения"],
    
    ["🔇 Акустические панели", 
     "• Многослойная структура\n• Огнестойкие покрытия\n• Архитектурная интеграция"],
    
    ["🚆 Транспортные решения", 
     "• Электродвигатели\n• Маглев-технологии\n• Бесшумные покрытия"]
]


# ======== СЛАЙД 8: Градостроительство ========
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_title(slide, "Урбанистические подходы")

# Схема с подписями


# ======== СЛАЙД 9: Успешные кейсы ========
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_title(slide, "Мировой опыт")

content = [
    ["🇯🇵 Япония", 
     "• Акустические тоннели\n• Вибрационные демпферы\n• Ночные ограничения"],
    
    ["🇩🇪 Германия", 
     "• Зоны тишины\n• Зеленые коридоры\n• «Тихий асфальт»"],
    
    ["🇸🇬 Сингапур", 
     
     "• Умные фасады\n• Подземные магистрали\n• Датчики шума"]
]



# ======== СЛАЙД 10: Заключение ========
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_title(slide, "Пути решения")

content = [[
    "✅ Технологические инновации\n"
    "✅ Грамотное зонирование\n"
    "✅ Экологическое образование\n"
    "✅ Международное сотрудничество\n\n"
    "«Тишина — основа гармоничного развития»"
]]


prs.save('noise_pollution_final.pptx')
print("Презентация успешно создана!")